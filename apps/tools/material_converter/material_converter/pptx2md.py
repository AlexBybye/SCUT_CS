"""PPTX -> Markdown per-slide extraction (faithful, no invention)."""
from __future__ import annotations

import re
from xml.etree import ElementTree as ET


def convert_pptx(path: str):
    """Return (slides, images, stats) where slides = [(slide_no, md_lines)], images=[(name, bytes)]."""
    from pptx import Presentation
    from io import BytesIO

    prs = Presentation(path)
    slides = []
    images = []
    stats = {"slides": len(prs.slides), "images": 0, "tables": 0, "notes": 0}
    img_idx = 0

    for idx, slide in enumerate(prs.slides, 1):
        lines = []
        # title placeholder first
        title = None
        try:
            if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
                t = slide.shapes.title.text.strip()
                if t:
                    title = t
        except Exception:
            title = None
        body_items = []

        def walk_shapes(shapes):
            nonlocal img_idx
            for shp in shapes:
                st = str(shp.shape_type)
                if getattr(shp, "shape_type", None) is not None and "GROUP" in st:
                    walk_shapes(getattr(shp, "shapes", []))
                    continue
                if shp.has_text_frame:
                    for para in shp.text_frame.paragraphs:
                        t = "".join(run.text for run in para.runs).strip() or (para.text or "").strip()
                        if not t:
                            continue
                        lvl = para.level or 0
                        is_title_shape = False
                        try:
                            if shp == slide.shapes.title:
                                is_title_shape = True
                        except Exception:
                            pass
                        body_items.append((lvl, t, id(shp), is_title_shape))
                elif "PICTURE" in st:
                    try:
                        image = shp.image
                        img_idx += 1
                        ext = image.ext
                        name = f"image-{img_idx:03d}.{ext}"
                        images.append((name, image.blob))
                        stats["images"] += 1
                        body_items.append((-1, f"![image](assets/{{ASSETS_DIR}}/{name})", id(shp), False))
                    except Exception:
                        pass
                elif getattr(shp, "has_table", False):
                    stats["tables"] += 1
                    tbl = shp.table
                    rows = []
                    for r in tbl.rows:
                        cells = [c.text.replace("|", "\\|").replace("\n", "<br>").strip() for c in r.cells]
                        rows.append(cells)
                    if rows:
                        ncols = max(len(r) for r in rows)
                        lines_out = ["", "| " + " | ".join(rows[0] + [""] * (ncols - len(rows[0]))) + " |",
                                     "|" + "---|" * ncols]
                        for r in rows[1:]:
                            lines_out.append("| " + " | ".join(r + [""] * (ncols - len(r))) + " |")
                        lines_out.append("")
                        body_items.append((-2, "\n".join(lines_out), id(shp), False))

        walk_shapes(slide.shapes)

        slide_md = []
        if title:
            slide_md.append(f"## {title}")
            slide_md.append("")
        seen_tables = set()
        for lvl, t, sid, is_title in body_items:
            if lvl == -2:
                if sid in seen_tables:
                    continue
                seen_tables.add(sid)
                slide_md.append(t)
            elif lvl == -1:
                slide_md.append(t)
            else:
                if is_title:
                    continue
                indent = "  " * lvl
                bullet = "- "
                slide_md.append(f"{indent}{bullet}{t}")
        notes = ""
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
                nt = slide.notes_slide.notes_text_frame.text.strip()
                if nt:
                    notes = nt
                    stats["notes"] += 1
        except Exception:
            pass
        slides.append((idx, slide_md, notes))
    return slides, images, stats
